#!/usr/bin/env python3
"""PPTX -> extraction.json parser (Phase 1).

Usage:
    python parser.py <pptx_path> <output_dir>

Contract:
    stdout    absolute path to the written extraction.json (and nothing else)
    stderr    warnings / errors
    exit 0    success, exit 1 failure
"""

import json
import os
import sys
from datetime import datetime, timezone

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.oxml.ns import qn

from extractors.colors import ThemeResolver
from extractors.shapes import build_shape, get_shape_type, emu_to_pt
from extractors.connectors import is_connector, extract_connector
from extractors.images import extract_image


def warn(msg):
    sys.stderr.write("[warn] %s\n" % msg)


def _is_group(shape):
    try:
        return shape.shape_type == MSO_SHAPE_TYPE.GROUP
    except Exception:
        return False


def _identity(left, top, width, height):
    return (left or 0, top or 0, width or 0, height or 0)


def _group_transform(group, parent_tf):
    """Return a function mapping a group child's EMU box to absolute slide EMU."""
    try:
        ax, ay, aw, ah = parent_tf(group.left, group.top, group.width, group.height)
        grp_pr = group._element.find(qn("p:grpSpPr"))
        xfrm = grp_pr.find(qn("a:xfrm"))
        ch_off = xfrm.find(qn("a:chOff"))
        ch_ext = xfrm.find(qn("a:chExt"))
        chx, chy = int(ch_off.get("x")), int(ch_off.get("y"))
        chcx, chcy = int(ch_ext.get("cx")), int(ch_ext.get("cy"))
        sx = aw / chcx if chcx else 1.0
        sy = ah / chcy if chcy else 1.0

        def tf(left, top, width, height):
            left, top = left or 0, top or 0
            width, height = width or 0, height or 0
            return (ax + (left - chx) * sx, ay + (top - chy) * sy, width * sx, height * sy)

        return tf
    except Exception as exc:
        warn("group transform failed for '%s'; using raw coords: %s" % (getattr(group, "name", "?"), exc))
        return parent_tf


def _flatten(shapes, tf):
    """Flatten groups recursively -> list of (shape, absolute_box_emu)."""
    result = []
    for shape in shapes:
        if _is_group(shape):
            result.extend(_flatten(shape.shapes, _group_transform(shape, tf)))
        else:
            result.append((shape, tf(shape.left, shape.top, shape.width, shape.height)))
    return result


def parse_slide(slide, index, output_dir, theme):
    flat = _flatten(slide.shapes, _identity)

    # pptx shape id -> "shape_<id>" key, for connector cross-referencing.
    shape_id_to_key = {
        sh.shape_id: "shape_%s" % sh.shape_id
        for sh, _ in flat
        if not is_connector(sh)
    }

    shapes_out, connectors_out, images_out = [], [], []

    for sh, box in flat:
        try:
            if is_connector(sh):
                connectors_out.append(extract_connector(sh, box, shape_id_to_key, theme))
            elif get_shape_type(sh) == "PICTURE":
                img = extract_image(sh, box, output_dir)
                if img:
                    images_out.append(img)
                else:
                    warn("slide %d: skipped unreadable image (shape %s)" % (index + 1, getattr(sh, "shape_id", "?")))
            else:
                shapes_out.append(build_shape(sh, box, theme))
        except Exception as exc:
            warn("slide %d: skipped shape %s: %s" % (index + 1, getattr(sh, "shape_id", "?"), exc))

    return {
        "slide_index": index,
        "slide_number": index + 1,
        "shapes": shapes_out,
        "connectors": connectors_out,
        "images": images_out,
    }


def parse(pptx_path, output_dir):
    prs = Presentation(pptx_path)
    theme = ThemeResolver(prs)

    sw = prs.slide_width
    sh = prs.slide_height

    slides = [parse_slide(slide, i, output_dir, theme) for i, slide in enumerate(prs.slides)]

    extraction = {
        "metadata": {
            "source_file": os.path.basename(pptx_path),
            "slide_width_emu": sw,
            "slide_height_emu": sh,
            "slide_width_pt": emu_to_pt(sw),
            "slide_height_pt": emu_to_pt(sh),
            "slide_count": len(slides),
            "extracted_at": datetime.now(timezone.utc).isoformat().replace("+00:00", "Z"),
        },
        "slides": slides,
    }

    os.makedirs(output_dir, exist_ok=True)
    json_path = os.path.join(output_dir, "extraction.json")
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(extraction, f, indent=2, ensure_ascii=False)

    return os.path.abspath(json_path)


def main(argv):
    if len(argv) != 3:
        sys.stderr.write("Usage: python parser.py <pptx_path> <output_dir>\n")
        return 1

    pptx_path, output_dir = argv[1], argv[2]

    if not os.path.isfile(pptx_path):
        sys.stderr.write("Error: file not found: %s\n" % pptx_path)
        return 1
    if not pptx_path.lower().endswith((".pptx", ".pptm")):
        sys.stderr.write("Error: expected a .pptx file: %s\n" % pptx_path)
        return 1

    try:
        json_path = parse(pptx_path, output_dir)
    except Exception as exc:
        sys.stderr.write("Error: parse failed: %s\n" % exc)
        return 1

    # Contract: the ONLY thing on stdout is the JSON path.
    sys.stdout.write(json_path + "\n")
    return 0


if __name__ == "__main__":
    sys.exit(main(sys.argv))
