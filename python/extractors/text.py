"""Text frame -> runs + Miro-safe HTML extraction.

Walks each paragraph's child elements in document order so that line breaks
(`<a:br/>`, which python-pptx surfaces as a vertical tab and otherwise drops
between runs) are preserved as `<br>`.
"""

import html as html_lib

from pptx.oxml.ns import qn

_ALIGN = {
    "LEFT": "left",
    "CENTER": "center",
    "RIGHT": "right",
    "JUSTIFY": "left",
    "JUSTIFY_LOW": "left",
    "DISTRIBUTE": "left",
    "THAI_DISTRIBUTE": "left",
}

_VANCHOR = {"TOP": "top", "MIDDLE": "middle", "BOTTOM": "bottom"}


def _enum_name(value):
    if value is None:
        return None
    return str(value).split(".")[-1].split(" ")[0]


def _align_name(alignment):
    return _ALIGN.get(_enum_name(alignment), "left")


def _vanchor_name(anchor):
    return _VANCHOR.get(_enum_name(anchor), "middle")


def _para_default_size(para):
    """Paragraph-level default run size (pt), or None."""
    pPr = para._pPr
    if pPr is not None:
        defRPr = pPr.find(qn("a:defRPr"))
        if defRPr is not None and defRPr.get("sz"):
            try:
                return int(defRPr.get("sz")) / 100.0
            except (TypeError, ValueError):
                return None
    return None


def _run_data(run, theme, fallback_size):
    font = run.font
    try:
        color = theme.color(font.color)
    except Exception:
        color = None
    try:
        size = float(font.size.pt) if font.size is not None else None
    except Exception:
        size = None
    if size is None:
        size = fallback_size
    return {
        "text": run.text,
        "font_name": font.name,
        "font_size_pt": size,
        "bold": bool(font.bold),
        "italic": bool(font.italic),
        "underline": bool(font.underline),
        "color": color,
    }


def _run_html(rd):
    esc = html_lib.escape(rd["text"])
    if rd["bold"]:
        esc = "<strong>%s</strong>" % esc
    if rd["italic"]:
        esc = "<em>%s</em>" % esc
    if rd["underline"]:
        esc = "<u>%s</u>" % esc
    return esc


def extract_text(text_frame, theme):
    """Return a dict with content/html/alignment/wrap/runs, or None when empty."""
    if text_frame is None:
        return None

    raw = text_frame.text
    if not raw.strip():
        return None

    runs_out = []
    html_parts = []
    first_align = "left"

    for i, para in enumerate(text_frame.paragraphs):
        align = _align_name(para.alignment)
        if i == 0:
            first_align = align

        para_default = _para_default_size(para)
        para_html = []
        run_iter = iter(para.runs)

        for child in para._p:
            tag = child.tag.split("}")[-1]
            if tag == "r":
                run = next(run_iter, None)
                if run is None:
                    continue
                rd = _run_data(run, theme, para_default)
                runs_out.append(rd)
                para_html.append(_run_html(rd))
            elif tag == "br":
                para_html.append("<br>")
            elif tag == "fld":
                # Field (e.g. slide number): keep its rendered text, no formatting.
                t = child.find(qn("a:t"))
                if t is not None and t.text:
                    para_html.append(html_lib.escape(t.text))

        html_parts.append("<p>%s</p>" % "".join(para_html))

    return {
        "content": raw.replace("\x0b", "\n").strip(),
        "html": "".join(html_parts),
        "alignment": first_align,
        "vertical_alignment": _vanchor_name(getattr(text_frame, "vertical_anchor", None)),
        # True = wrap (square), False = no-wrap (text overflows box), None = inherit.
        "wrap": text_frame.word_wrap,
        "runs": runs_out,
    }
