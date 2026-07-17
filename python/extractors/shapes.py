"""Shape geometry, type, fill, and border extraction."""

from pptx.enum.shapes import MSO_SHAPE_TYPE
from .text import extract_text

EMU_PER_PT = 12700

# PPTX auto-shape type -> normalized type consumed by the Node mapper.
AUTOSHAPE_TYPE_MAP = {
    "RECTANGLE": "RECTANGLE",
    "ROUNDED_RECTANGLE": "ROUNDED_RECTANGLE",
    "ROUND_1_RECTANGLE": "ROUNDED_RECTANGLE",
    "OVAL": "ELLIPSE",
    "DIAMOND": "DIAMOND",
    "ISOCELES_TRIANGLE": "TRIANGLE",
    "TRIANGLE": "TRIANGLE",
    "RIGHT_TRIANGLE": "TRIANGLE",
    "PARALLELOGRAM": "PARALLELOGRAM",
    "TRAPEZOID": "TRAPEZOID",
    "HEXAGON": "HEXAGON",
    "OCTAGON": "OCTAGON",
    "PENTAGON": "PENTAGON",
    "CROSS": "CROSS",
    "STAR_5_POINT": "STAR",
    "STAR_4_POINT": "STAR",
    "STAR_6_POINT": "STAR",
    "CLOUD": "CLOUD",
    "CAN": "CYLINDER",
    "CYLINDER": "CYLINDER",
    "CHEVRON": "CHEVRON",
    "RIGHT_ARROW": "RIGHT_ARROW",
}

_DASH_MAP = {
    "SOLID": "SOLID",
    "DASH": "DASHED",
    "LONG_DASH": "DASHED",
    "DASH_DOT": "DASHED",
    "LONG_DASH_DOT": "DASHED",
    "LONG_DASH_DOT_DOT": "DASHED",
    "ROUND_DOT": "DOTTED",
    "SQUARE_DOT": "DOTTED",
    "DOT": "DOTTED",
}


def emu_to_pt(value):
    if value is None:
        return 0.0
    return round(value / EMU_PER_PT, 4)


def _enum_name(value):
    if value is None:
        return None
    return str(value).split(".")[-1].split(" ")[0]


def get_shape_type(shape):
    st = shape.shape_type
    if st == MSO_SHAPE_TYPE.TEXT_BOX:
        return "TEXT_BOX"
    if st == MSO_SHAPE_TYPE.PICTURE:
        return "PICTURE"
    if st == MSO_SHAPE_TYPE.LINE:
        return "LINE"
    if st == MSO_SHAPE_TYPE.PLACEHOLDER:
        return "TEXT_BOX"
    if st == MSO_SHAPE_TYPE.AUTO_SHAPE:
        try:
            auto_type = _enum_name(shape.auto_shape_type)
            return AUTOSHAPE_TYPE_MAP.get(auto_type, "RECTANGLE")
        except Exception:
            return "RECTANGLE"
    return "RECTANGLE"


def extract_fill(fill, theme):
    try:
        ftype = fill.type
    except Exception:
        return {"type": "NONE", "color": None}

    if ftype is None:
        return {"type": "NONE", "color": None}

    tname = _enum_name(ftype)

    if tname == "SOLID":
        try:
            return {"type": "SOLID", "color": theme.color(fill.fore_color)}
        except Exception:
            return {"type": "SOLID", "color": None}

    if tname == "GRADIENT":
        # Miro has no gradients; approximate with the first gradient stop.
        try:
            stops = list(fill.gradient_stops)
            color = theme.color(stops[0].color) if stops else None
            return {"type": "GRADIENT", "color": color}
        except Exception:
            return {"type": "GRADIENT", "color": None}

    return {"type": tname, "color": None}


def extract_line(line, theme):
    try:
        color = theme.color(line.color)
    except Exception:
        color = None

    try:
        width_pt = round(line.width / EMU_PER_PT, 4) if line.width else None
    except Exception:
        width_pt = None

    try:
        style = _DASH_MAP.get(_enum_name(line.dash_style), "SOLID")
    except Exception:
        style = "SOLID"

    return {"color": color, "width_pt": width_pt, "style": style}


def _safe_fill(shape, theme):
    try:
        return extract_fill(shape.fill, theme)
    except Exception:
        return None


def _safe_border(shape, theme):
    try:
        return extract_line(shape.line, theme)
    except Exception:
        return None


def build_shape(shape, box, theme):
    """box = (left, top, width, height) in EMU, already group-transformed."""
    left, top, width, height = box
    text = None
    if getattr(shape, "has_text_frame", False):
        try:
            text = extract_text(shape.text_frame, theme)
        except Exception:
            text = None

    return {
        "id": "shape_%s" % shape.shape_id,
        "pptx_shape_id": shape.shape_id,
        "name": shape.name,
        "type": get_shape_type(shape),
        "x_pt": emu_to_pt(left),
        "y_pt": emu_to_pt(top),
        "width_pt": emu_to_pt(width),
        "height_pt": emu_to_pt(height),
        "rotation": float(shape.rotation or 0.0),
        "fill": _safe_fill(shape, theme),
        "border": _safe_border(shape, theme),
        "text": text,
    }
